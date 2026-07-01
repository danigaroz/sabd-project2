"""Generate the SABD Project 2 slide deck (15-min oral presentation).

Uses python-pptx, no external dependencies beyond the library itself.

Run:
    pip install python-pptx
    python make_slides.py

Output:
    sabd_p2_slides.pptx (in the same directory as this script)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ----- palette --------------------------------------------------------------
NAVY    = RGBColor(0x1E, 0x27, 0x61)   # primary dark
ICE     = RGBColor(0xCA, 0xDC, 0xFC)   # soft secondary
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREY    = RGBColor(0x33, 0x33, 0x33)
LGREY   = RGBColor(0x66, 0x66, 0x66)
ORANGE  = RGBColor(0xFF, 0x9F, 0x1C)   # accent: batch path / Q2
GREEN   = RGBColor(0x4C, 0xAF, 0x50)   # success / OK
RED     = RGBColor(0xE5, 0x53, 0x4F)   # failure / OOM

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ----- helpers --------------------------------------------------------------

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=GREY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_rich_lines(slide, x, y, w, h, lines, *, size=16, color=GREY,
                   font="Calibri"):
    """lines = list of (text, is_bold, color_or_None)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    for i, (txt, bold, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = txt
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = col if col is not None else color
        p.space_after = Pt(4)
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=16, color=GREY,
                bullet="•", line_gap=8):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    for i, txt in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"{bullet}  {txt}"
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.color.rgb = color
        p.space_after = Pt(line_gap)
    return tb


def add_rect(slide, x, y, w, h, *, fill=ICE, line=None, line_w=0.75,
             rounded=True, shadow=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is not None:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    else:
        sh.line.fill.background()
    if not shadow:
        sh.shadow.inherit = False
    return sh


def add_card(slide, x, y, w, h, title, body_lines, *, head_color=NAVY,
             card_fill=WHITE, head_size=18, body_size=14):
    add_rect(slide, x, y, w, h, fill=card_fill, line=ICE, line_w=0.75)
    pad = Inches(0.2)
    add_text(slide, x + pad, y + pad, w - 2 * pad, Inches(0.4),
             title, size=head_size, bold=True, color=head_color)
    add_bullets(slide, x + pad, y + Inches(0.65), w - 2 * pad, h - Inches(0.85),
                body_lines, size=body_size, color=GREY, line_gap=4)


def slide_title_bar(slide, title, subtitle=None):
    """Header band: title + thin subtitle. No accent lines."""
    add_text(slide, Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.7),
             title, size=30, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.55), Inches(1.0), Inches(12.2), Inches(0.4),
                 subtitle, size=15, color=LGREY)


def slide_page_number(slide, n, total):
    add_text(slide, Inches(12.0), Inches(7.1), Inches(1.0), Inches(0.3),
             f"{n} / {total}", size=10, color=LGREY, align=PP_ALIGN.RIGHT)


# ----- presentation ---------------------------------------------------------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

TOTAL = 12

# ============================================================================
# Slide 1 — Title
# ============================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, NAVY)

# Big title centred vertically
add_text(s, Inches(0.7), Inches(1.5), Inches(12.0), Inches(0.5),
         "SABD — Project 2", size=18, bold=True, color=ICE,
         align=PP_ALIGN.LEFT)
add_text(s, Inches(0.7), Inches(2.1), Inches(12.0), Inches(1.6),
         "Real-Time Analysis of US Flight Delays\nwith Apache Flink",
         size=42, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s, Inches(0.7), Inches(3.9), Inches(12.0), Inches(0.6),
         "Event-Time Windowing, Stream Simulation,",
         size=20, color=ICE, align=PP_ALIGN.LEFT)
add_text(s, Inches(0.7), Inches(4.3), Inches(12.0), Inches(0.6),
         "and a Hybrid Stream/Batch Fallback for Top-K Ranking",
         size=20, color=ICE, align=PP_ALIGN.LEFT)

# Author block
add_text(s, Inches(0.7), Inches(5.8), Inches(12.0), Inches(0.45),
         "Daniel Garoz", size=20, bold=True, color=WHITE,
         align=PP_ALIGN.LEFT)
add_text(s, Inches(0.7), Inches(6.2), Inches(12.0), Inches(0.4),
         "Università degli Studi di Roma Tor Vergata · A.A. 2025/26",
         size=14, color=ICE, align=PP_ALIGN.LEFT)

# ============================================================================
# Slide 2 — Project goals
# ============================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
slide_title_bar(s, "Project goals",
                "Single-student team — Q1 + Q2 mandatory, Q3 optional (skipped)")

# Two columns: Q1 and Q2
add_card(s, Inches(0.7), Inches(1.8), Inches(5.8), Inches(4.8),
         "Q1 — Hourly carrier statistics",
         [
             "Carriers AA, DL, UA, WN",
             "Tumbling event-time windows of 1 hour",
             "Per (window, carrier): total/completed/cancelled/diverted",
             "Mean DEP_DELAY of non-cancelled",
             "Cancellation rate, late-departure rate",
             "Output: CSV with 10 columns",
         ],
         head_color=NAVY, head_size=20, body_size=15)

add_card(s, Inches(6.85), Inches(1.8), Inches(5.8), Inches(4.8),
         "Q2 — Top-10 airports by severe delays",
         [
             "ALL carriers, group by ORIGIN_AIRPORT_ID",
             "Severe = non-cancelled & non-diverted & DEP_DELAY > 30 min",
             "3 event-time windows: 1h, 6h, full dataset",
             "Filter: ≥ 30 non-cancelled non-diverted flights",
             "Top-10 by descending severe-count",
             "Up to 20 most-delayed flights per airport",
         ],
         head_color=ORANGE, head_size=20, body_size=15)

# Dataset footer
add_text(s, Inches(0.7), Inches(6.85), Inches(12.0), Inches(0.4),
         "Dataset: US BTS — 2,229,453 flight events · 1 Jan – 30 Apr 2025",
         size=12, color=LGREY)
slide_page_number(s, 2, TOTAL)

# ============================================================================
# Slide 3 — System architecture
# ============================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
slide_title_bar(s, "System architecture",
                "3 Docker containers · bind-mounted project tree · custom PyFlink image")

# Diagram region: x in [0.7, 12.6], y in [1.8, 6.8]
def box(x, y, w, h, label, sub=None, *, fill=ICE, head_color=NAVY,
        head_size=14, sub_size=10):
    add_rect(s, x, y, w, h, fill=fill, line=NAVY, line_w=0.75)
    if sub:
        add_text(s, x, y + Inches(0.15), w, Inches(0.35),
                 label, size=head_size, bold=True, color=head_color,
                 align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(0.5), w, Inches(0.3),
                 sub, size=sub_size, color=LGREY,
                 align=PP_ALIGN.CENTER)
    else:
        add_text(s, x, y + Inches(0.25), w, Inches(0.5),
                 label, size=head_size, bold=True, color=head_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def arrow(x1, y1, x2, y2, dashed=False, color=NAVY):
    conn = s.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = straight
    conn.line.color.rgb = color
    conn.line.width = Pt(1.5)
    if dashed:
        # python-pptx: set dash via XML
        from pptx.oxml.ns import qn
        ln = conn.line._get_or_add_ln()
        # remove existing prstDash if any
        for d in ln.findall(qn('a:prstDash')):
            ln.remove(d)
        from lxml import etree
        prstDash = etree.SubElement(ln, qn('a:prstDash'))
        prstDash.set('val', 'dash')
    # end arrow head
    from pptx.oxml.ns import qn
    ln = conn.line._get_or_add_ln()
    from lxml import etree
    for t in ln.findall(qn('a:tailEnd')):
        ln.remove(t)
    tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
    tailEnd.set('type', 'triangle')
    tailEnd.set('w', 'med')
    tailEnd.set('h', 'med')

# Layout: events.csv at top center, branches left (streaming) and right (batch)
csv_x = Inches(5.4); csv_y = Inches(2.0); csv_w = Inches(2.5); csv_h = Inches(0.9)
box(csv_x, csv_y, csv_w, csv_h, "events.csv", "(pre-sorted by event_time)",
    fill=ICE, head_color=NAVY, head_size=15, sub_size=10)

# left branch (streaming): feeder, flink, q1
left_x = Inches(2.4); right_x = Inches(8.5)
feeder_y = Inches(3.6); flink_y = Inches(4.9); q1_y = Inches(6.2)
node_w = Inches(2.5); node_h = Inches(0.9)
box(left_x, feeder_y, node_w, node_h, "Feeder (Python)", "TCP :9999",
    fill=ICE, head_color=NAVY, head_size=14, sub_size=10)
box(left_x, flink_y, node_w, node_h, "Flink 1.20", "JobManager + TaskManager",
    fill=ICE, head_color=NAVY, head_size=14, sub_size=10)
box(left_x, q1_y, node_w, node_h, "Q1 CSV", "(Results/q1)",
    fill=ICE, head_color=NAVY, head_size=14, sub_size=10)

# right branch (batch): post-process, q2 csvs
box(right_x, feeder_y, node_w, node_h, "Post-process", "(pandas)",
    fill=RGBColor(0xFF, 0xE9, 0xCB), head_color=ORANGE, head_size=14, sub_size=10)
box(right_x, flink_y, node_w, node_h, "Q2 CSVs", "(Results/q2_1h, _6h, _global)",
    fill=RGBColor(0xFF, 0xE9, 0xCB), head_color=ORANGE, head_size=14, sub_size=10)

# Arrows: csv -> feeder, feeder -> flink, flink -> q1
arrow(csv_x + csv_w / 2, csv_y + csv_h,
      left_x + node_w / 2, feeder_y, color=NAVY)
arrow(left_x + node_w / 2, feeder_y + node_h,
      left_x + node_w / 2, flink_y, color=NAVY)
arrow(left_x + node_w / 2, flink_y + node_h,
      left_x + node_w / 2, q1_y, color=NAVY)

# csv -> post-process, post-process -> q2 csvs (dashed)
arrow(csv_x + csv_w / 2, csv_y + csv_h,
      right_x + node_w / 2, feeder_y, color=ORANGE, dashed=True)
arrow(right_x + node_w / 2, feeder_y + node_h,
      right_x + node_w / 2, flink_y, color=ORANGE, dashed=True)

# Legend
add_text(s, Inches(0.7), Inches(6.85), Inches(12.0), Inches(0.4),
         "Solid (navy) = streaming path (Q1)   |   Dashed (orange) = batch path (Q2)",
         size=11, color=LGREY)
slide_page_number(s, 3, TOTAL)

# ============================================================================
# Slide 4 — Stream simulation
# ============================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
slide_title_bar(s, "Stream simulation",
                "Replaying a historical dataset as a real-time event-time stream")

# 3 cards: pre-sort, feeder, watermark
card_w = Inches(4.0); card_h = Inches(4.6); card_y = Inches(1.8)

add_card(s, Inches(0.55), card_y, card_w, card_h,
         "1. Pre-sort once",
         [
             "Read 4 monthly CSVs (Jan–Apr 2025)",
             "Build event_time from YEAR, MONTH,",
             "  DAY_OF_MONTH, CRS_DEP_TIME",
             "Sort by event_time → events.csv",
             "Why: avoids out-of-order at source",
             "  → simplifies watermarking",
         ],
         head_color=NAVY, head_size=18, body_size=13)

add_card(s, Inches(4.7), card_y, card_w, card_h,
         "2. TCP feeder + acceleration",
         [
             "Python TCP server on port 9999",
             "Replays one line per emit",
             "Sleep Δt_wall = (et_{n+1}−et_n) / α",
             "α = 60,000 → 120 days replays in ~3 min",
             "α only compresses wall-clock,",
             "  event-time semantics unchanged",
         ],
         head_color=NAVY, head_size=18, body_size=13)

add_card(s, Inches(8.85), card_y, card_w, card_h,
         "3. Watermark strategy",
         [
             "for_bounded_out_of_orderness(60s)",
             "Events arrive ordered by construction;",
             "  60s margin = defensive jitter buffer",
             "Smallest window is 1h → 60s ≪ 1h",
             "Null DEP_DELAY treated as 0",
             "  (spec-permitted, branch-free aggregator)",
         ],
         head_color=NAVY, head_size=18, body_size=13)

slide_page_number(s, 4, TOTAL)

# ============================================================================
# Slide 5 — Q1 PyFlink pipeline
# ============================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
slide_title_bar(s, "Query 1 — DataStream pipeline",
                "Textbook Flink streaming: source → watermark → keyBy → window → sink")

# Vertical pipeline diagram on left, code-style description on right
steps = [
    "socketTextStream(feeder, 9999)",
    "flat_map(parse + filter ∈ {AA,DL,UA,WN})",
    "assign_timestamps_and_watermarks(60s)",
    "key_by(carrier)",
    "window(TumblingEventTimeWindows(1h))",
    "aggregate(Q1Aggregator, Q1WindowFn)",
    "FileSink (Results/q1)",
]
sx = Inches(0.6); sy = Inches(1.8)
sw = Inches(6.5); sh = Inches(0.55); gap = Inches(0.15)
for i, t in enumerate(steps):
    y = sy + (sh + gap) * i
    add_rect(s, sx, y, sw, sh, fill=ICE, line=NAVY, line_w=0.75)
    add_text(s, sx + Inches(0.15), y, sw - Inches(0.3), sh, t,
             size=14, color=NAVY, bold=False, anchor=MSO_ANCHOR.MIDDLE,
             font="Consolas")
    if i < len(steps) - 1:
        y2 = y + sh
        conn = s.shapes.add_connector(1, sx + sw / 2, y2,
                                       sx + sw / 2, y2 + gap)
        conn.line.color.rgb = NAVY
        conn.line.width = Pt(1.5)

# Right: aggregator narrative
nx = Inches(7.4); ny = Inches(1.8); nw = Inches(5.5); nh = Inches(5.0)
add_text(s, nx, ny, nw, Inches(0.5),
         "Aggregator state — 7 primitives:", size=18, bold=True, color=NAVY)
add_text(s, nx, ny + Inches(0.55), nw, Inches(1.6),
         "(num, completed, cancelled, diverted,\n"
         " sum_dep_delay, count_non_cancelled,\n"
         " late_non_cancelled)",
         size=13, color=GREY, font="Consolas")
add_text(s, nx, ny + Inches(2.5), nw, Inches(0.5),
         "Window function emits one CSV line per (window, carrier):",
         size=16, bold=True, color=NAVY)
add_bullets(s, nx, ny + Inches(3.0), nw, Inches(2.5),
            [
                "mean_delay = sum / count_non_cancelled",
                "cancellation_rate = 100 · cancelled / num",
                "late_dep_rate = 100 · late_non_cancelled / count_non_cancelled",
                "Constant-time updates → ~10,000 ev/s sustained",
            ], size=14)
slide_page_number(s, 5, TOTAL)

# ============================================================================
# Slide 6 — Q1 output sample
# ============================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
slide_title_bar(s, "Query 1 — Output",
                "9,765 (window, carrier) rows over 120 days of event time")

# Sample table
header = ["window_start", "window_end", "airline", "num", "compl", "canc",
          "div", "dep_delay_mean", "canc_rate", "late_rate"]
rows = [
    ["2025-01-01 00:00", "2025-01-01 01:00", "AA", "13", "13", "0", "0",
     "7.38", "0.00", "30.77"],
    ["2025-01-01 00:00", "2025-01-01 01:00", "DL", "1",  "1",  "0", "0",
     "-6.00", "0.00", "0.00"],
    ["2025-01-01 00:00", "2025-01-01 01:00", "UA", "7",  "7",  "0", "0",
     "10.57", "0.00", "14.29"],
    ["2025-04-20 08:00", "2025-04-20 09:00", "AA", "215","213","2","0",
     "7.52", "0.93", "18.31"],
    ["2025-04-20 08:00", "2025-04-20 09:00", "DL", "251","251","0","0",
     "1.06", "0.00", "8.37"],
]

# Header row
hx = Inches(0.55); hy = Inches(1.9)
col_w = [1.5, 1.5, 0.8, 0.8, 0.9, 0.9, 0.7, 1.5, 1.2, 1.2]  # inches
col_x = [hx]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + Inches(w))

# Header background bar
add_rect(s, hx, hy, sum(Inches(w) for w in col_w), Inches(0.45),
         fill=NAVY, rounded=False)
for i, (txt, w) in enumerate(zip(header, col_w)):
    add_text(s, col_x[i], hy, Inches(w), Inches(0.45), txt,
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

# Data rows
ry = hy + Inches(0.45)
for r_i, row in enumerate(rows):
    bgfill = WHITE if r_i % 2 == 0 else RGBColor(0xF4, 0xF6, 0xFC)
    add_rect(s, hx, ry, sum(Inches(w) for w in col_w), Inches(0.4),
             fill=bgfill, rounded=False)
    for i, (val, w) in enumerate(zip(row, col_w)):
        add_text(s, col_x[i], ry, Inches(w), Inches(0.4), val,
                 size=11, color=GREY, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, font="Consolas")
    ry += Inches(0.4)

# Insight callout
add_text(s, Inches(0.55), Inches(5.4), Inches(12.0), Inches(0.5),
         "Sanity check:", size=16, bold=True, color=NAVY)
add_bullets(s, Inches(0.55), Inches(5.9), Inches(12.0), Inches(1.3),
            [
                "AA at 2025-01-01 00:00–01:00: 13 flights, all completed, mean delay 7.38 min, 30.77% late (>15 min)",
                "AA at 2025-04-20 08:00–09:00: 215 flights, 2 cancelled (0.93%), 213 completed → busy hub-style traffic",
                "Time-of-day pattern (early morning low traffic, late morning peak) matches BTS dataset characteristics",
            ], size=13)
slide_page_number(s, 6, TOTAL)

# ============================================================================
# Slide 7 — Q2 four PyFlink iterations
# ============================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
slide_title_bar(s, "Query 2 — Four PyFlink iterations",
                "Each version exposed a different failure mode; honest engineering record")

# 2x2 grid of cards
cw = Inches(6.1); ch = Inches(2.5)
gx = Inches(0.55); gy = Inches(1.7); gap_x = Inches(0.15); gap_y = Inches(0.18)

versions = [
    ("v1 — ProcessAllWindowFunction",
     ["window_all + Python top-10 in process()",
      "Failure: OOM at ~900K events",
      "Cause: 120-day window accumulated all events",
      "  in Beam worker heap (~134 MB)"],
     RED),
    ("v2 — keyBy + PICKLED accumulator",
     ["Pre-agg keyed by airport_id",
      "Acc: (num,severe,sum,max,top20_list)",
      "Throughput: ~500 ev/s (2–3 orders worse than Q1)",
      "Cause: per-event pickle/unpickle dominates"],
     ORANGE),
    ("v3 — v2 + parallelism = 4",
     ["Throughput regressed to ~300 ev/s",
      "Source is socket → fixed parallelism 1",
      "Adding -p 4 only adds Beam coordination",
      "  overhead without source parallelism"],
     ORANGE),
    ("v4 — primitive acc + bundle tuning",
     ["Removed top-20 list, primitive TUPLE acc",
      "bundle.size = 10,000 (default 1000)",
      "~540 ev/s; failed at 68% (1.53M / 2.23M)",
      "TaskManager exhausted 1.7 GB process budget"],
     ORANGE),
]

for i, (title, body, head_col) in enumerate(versions):
    col = i % 2; row = i // 2
    x = gx + col * (cw + gap_x)
    y = gy + row * (ch + gap_y)
    add_card(s, x, y, cw, ch, title, body,
             head_color=head_col, head_size=16, body_size=12)

slide_page_number(s, 7, TOTAL)

# ============================================================================
# Slide 8 — Q2 hybrid stream/batch design
# ============================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
slide_title_bar(s, "Query 2 — Final design: hybrid stream/batch",
                "Spec allows 1-student teams to «simulare lo stream e scrivere risultati CSV»")

# Two columns: rationale + numbers
add_card(s, Inches(0.55), Inches(1.8), Inches(6.2), Inches(4.8),
         "Decision: pandas batch over events.csv",
         [
             "Reads same events.csv that drives the feeder",
             "Groups by (event-time window, airport_id)",
             "Pandas floor(\"1h\") and floor(\"6h\") for buckets",
             "Global window = single bucket over full range",
             "Reconstructs top-20 severely-delayed list",
             "All output CSVs match spec schema exactly",
         ],
         head_color=ORANGE, head_size=18, body_size=14)

add_card(s, Inches(7.0), Inches(1.8), Inches(5.8), Inches(4.8),
         "Why it's defensible",
         [
             "Event-time semantics preserved (windows on event_time, not wall-clock)",
             "Spec quote: «simulare lo stream di input»",
             "  Feeder = stream simulation",
             "  Batch processor = downstream consumer",
             "Q1 demonstrates streaming mastery (Flink DataStream API)",
             "Q2 turnaround: hours → 30 s",
             "Same output schema, identical results",
         ],
         head_color=NAVY, head_size=18, body_size=14)

slide_page_number(s, 8, TOTAL)

# ============================================================================
# Slide 9 — Benchmarks Q1
# ============================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
slide_title_bar(s, "Benchmarks — Q1 throughput vs. parallelism",
                "α = 60,000 · 3 runs · output identical (verified with diff)")

# Centered table
header = ["parallelism", "total time (s)", "throughput (ev/s)", "output rows"]
rows = [
    ["1", "222.9", "10,002", "9,765"],
    ["2", "222.9", "10,002", "9,765"],
    ["4", "222.9", "10,002", "9,765"],
]
col_w = [2.0, 2.5, 2.8, 2.3]
hx = Inches((13.333 - sum(col_w)) / 2); hy = Inches(2.0)
col_x = [hx]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + Inches(w))

add_rect(s, hx, hy, sum(Inches(w) for w in col_w), Inches(0.5),
         fill=NAVY, rounded=False)
for i, (txt, w) in enumerate(zip(header, col_w)):
    add_text(s, col_x[i], hy, Inches(w), Inches(0.5), txt,
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

ry = hy + Inches(0.5)
for r_i, row in enumerate(rows):
    bgfill = WHITE if r_i % 2 == 0 else RGBColor(0xF4, 0xF6, 0xFC)
    add_rect(s, hx, ry, sum(Inches(w) for w in col_w), Inches(0.5),
             fill=bgfill, rounded=False)
    for i, (val, w) in enumerate(zip(row, col_w)):
        add_text(s, col_x[i], ry, Inches(w), Inches(0.5), val,
                 size=13, color=GREY, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, font="Consolas")
    ry += Inches(0.5)

# Observation block
add_text(s, Inches(0.55), Inches(4.5), Inches(12.0), Inches(0.5),
         "Why is the table flat?", size=18, bold=True, color=NAVY)
add_bullets(s, Inches(0.55), Inches(5.0), Inches(12.0), Inches(2.0),
            [
                "socketTextStream has intrinsic parallelism 1 — single Java reader thread consumes the TCP socket",
                "Downstream keyBy is on 4 carriers only — p∈{2,4} subtasks cannot shard a single key group",
                "Take-away: to scale, need (i) parallelisable source (Kafka) + (ii) higher key cardinality",
            ], size=14)
slide_page_number(s, 9, TOTAL)

# ============================================================================
# Slide 10 — Q2 throughput across iterations
# ============================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
slide_title_bar(s, "Q2 — PyFlink iterations: throughput collapse",
                "Documenting the design pressure that led to the batch fallback")

header = ["version", "ev/s", "outcome"]
rows = [
    ["v1 — ProcessAllWindowFn", "—",      "OOM at ~900 K events"],
    ["v2 — PICKLED acc",         "500",    "Throughput-bound (~73 min ETA)"],
    ["v3 — v2 + parallelism = 4", "300",    "Regression vs v2"],
    ["v4 — primitive acc + bundle 10k", "540", "Failed at 1.53M / 2.23M (68%)"],
    ["FINAL — pandas batch",     "~75,000", "Full output in ~30 s"],
]
col_w = [4.5, 1.6, 5.4]
hx = Inches((13.333 - sum(col_w)) / 2); hy = Inches(2.0)
col_x = [hx]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + Inches(w))

add_rect(s, hx, hy, sum(Inches(w) for w in col_w), Inches(0.5),
         fill=NAVY, rounded=False)
for i, (txt, w) in enumerate(zip(header, col_w)):
    add_text(s, col_x[i], hy, Inches(w), Inches(0.5), txt,
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

ry = hy + Inches(0.5)
for r_i, row in enumerate(rows):
    is_final = r_i == len(rows) - 1
    bgfill = (RGBColor(0xE8, 0xF5, 0xE9) if is_final
              else (WHITE if r_i % 2 == 0 else RGBColor(0xF4, 0xF6, 0xFC)))
    add_rect(s, hx, ry, sum(Inches(w) for w in col_w), Inches(0.5),
             fill=bgfill, rounded=False)
    for i, (val, w) in enumerate(zip(row, col_w)):
        col = GREEN if is_final else GREY
        bold = is_final
        add_text(s, col_x[i], ry, Inches(w), Inches(0.5), val,
                 size=13, color=col, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, font="Consolas", bold=bold)
    ry += Inches(0.5)

# Key insight
add_text(s, Inches(0.55), Inches(5.5), Inches(12.0), Inches(0.5),
         "Root cause", size=18, bold=True, color=NAVY)
add_text(s, Inches(0.55), Inches(6.0), Inches(12.0), Inches(1.2),
         "Beam-based PyFlink with ~360 keyed states + 2–3 windows saturates Java↔Python RPC. "
         "Pure Java/Scala Flink would not exhibit this; porting was out of scope for a "
         "1-student team. The hybrid design ships the spec result reliably.",
         size=14, color=GREY)
slide_page_number(s, 10, TOTAL)

# ============================================================================
# Slide 11 — Lessons learned
# ============================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
slide_title_bar(s, "Lessons learned",
                "Three engineering take-aways for streaming on a single-host budget")

cards = [
    ("Source parallelism matters more than job -p",
     [
         "socketTextStream is paralelism = 1",
         "Adding more subtasks cannot beat the source",
         "Kafka/Pulsar partitions = real source parallelism",
     ],
     NAVY),
    ("PyFlink ≠ Flink for hot paths",
     [
         "Beam RPC dominates when accumulators",
         "  are complex or keys are many",
         "Primitive TUPLE accumulators help a lot",
         "Bundle size tuning amortises RPC cost",
     ],
     NAVY),
    ("Spec relaxation is a feature, not a defeat",
     [
         "Spec allows batch for 1-student teams",
         "Designing for that relaxation = correct call",
         "Q1 proves streaming chops; Q2 ships results",
     ],
     ORANGE),
]
cw = Inches(4.1); ch = Inches(4.3); cy = Inches(1.8)
for i, (title, body, col) in enumerate(cards):
    x = Inches(0.55) + i * (cw + Inches(0.15))
    add_card(s, x, cy, cw, ch, title, body,
             head_color=col, head_size=16, body_size=13)

slide_page_number(s, 11, TOTAL)

# ============================================================================
# Slide 12 — Closing / Q&A
# ============================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, NAVY)

add_text(s, Inches(0.7), Inches(2.0), Inches(12.0), Inches(0.8),
         "Thank you — questions?", size=46, bold=True, color=WHITE,
         align=PP_ALIGN.LEFT)
add_text(s, Inches(0.7), Inches(3.2), Inches(12.0), Inches(0.6),
         "Daniel Garoz", size=22, bold=True, color=ICE,
         align=PP_ALIGN.LEFT)
add_text(s, Inches(0.7), Inches(3.7), Inches(12.0), Inches(0.4),
         "danigarochini@gmail.com", size=16, color=ICE, align=PP_ALIGN.LEFT)

# Bottom callout: deliverables
add_text(s, Inches(0.7), Inches(5.4), Inches(12.0), Inches(0.5),
         "Deliverables in the repository:", size=15, bold=True, color=ICE)
add_bullets(s, Inches(0.7), Inches(5.9), Inches(12.0), Inches(1.4),
            [
                "src/   — query1.py (PyFlink), query2_batch.py (pandas), feeder.py",
                "Results/   — q1 (9,765 rows), q2_{1h,6h,global}_final.csv",
                "Report/   — report.pdf (IEEE) + ai_declaration.pdf",
            ], size=14, color=ICE)

# ----- write & exit ---------------------------------------------------------
import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "sabd_p2_slides.pptx")
prs.save(out)
print(f"Wrote {out}")
